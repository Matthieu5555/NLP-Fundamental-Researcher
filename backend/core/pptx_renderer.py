"""
PowerPoint slide renderer.

Consumes the context dict produced by slides_generator.build_slides_context()
and renders an editable .pptx deck using python-pptx. All text, tables, and
metric cards are native PowerPoint objects so analysts can edit them directly.

Single entry point: generate_pptx(context) -> bytes
"""

import io
import logging
from typing import Any, Dict, List, Optional, Sequence

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from backend.core.pptx_template import (
    CONTENT_HEIGHT,
    CONTENT_LEFT,
    CONTENT_TOP,
    CONTENT_WIDTH,
    MARGIN_BOTTOM,
    MARGIN_LEFT,
    MARGIN_RIGHT,
    SLIDE_HEIGHT,
    SLIDE_WIDTH,
    FontStyle,
    PptxColors,
    add_brand_stripe,
    add_bullet_list,
    add_footer,
    add_metric_card,
    add_paragraph,
    add_slide_title,
    apply_font,
    apply_gradient_fill,
    build_colors,
    create_table,
    font_styles,
    hex_to_pptx,
    set_cell,
    set_text,
    style_table,
)

logger = logging.getLogger(__name__)

# Number of slides before content starts (title + TOC)
_PREAMBLE_SLIDES = 2


def generate_pptx(context: Dict[str, Any]) -> bytes:
    """Generate an editable PowerPoint deck from a slides context dict.

    The context dict is the same one produced by build_slides_context() in
    slides_generator.py. This function is a pure transformation from that
    dict to PPTX bytes, with no LLM calls or data fetching.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    colors = build_colors(context.get("primary_color"))
    styles = font_styles(colors)

    slide_num = [0]  # mutable counter for closures

    def new_slide() -> Any:
        """Add a blank slide, apply brand stripe and footer, increment counter."""
        layout = prs.slide_layouts[6]  # blank layout
        slide = prs.slides.add_slide(layout)
        slide_num[0] += 1
        add_brand_stripe(slide, colors)
        add_footer(slide, context.get("firm_name", ""), slide_num[0], colors, styles)
        return slide

    # Build slides in order. Each builder returns early if data is missing.
    builders = [
        _slide_title,
        _slide_toc,
        _slide_exec_summary,
        _slide_thesis,
        _slide_conviction,
        _slide_financial_overview,
        _slide_football_field,
        _slide_dcf,
        _slide_scenarios,
        _slide_sensitivity,
        _slide_earnings,
        _slide_comps,
        _slide_precedents,
        _slide_catalysts_risks,
        _slide_monitoring,
        _slide_strategy,
        _slide_industry,
        _slide_moat,
        _slide_investment_cases,
        _slide_sources,
        _slide_disclaimer,
    ]

    for builder in builders:
        try:
            builder(new_slide, context, colors, styles)
        except Exception:
            logger.warning("Slide builder %s failed, skipping", builder.__name__, exc_info=True)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# =============================================================================
# INDIVIDUAL SLIDE BUILDERS
#
# Each takes (new_slide_fn, context, colors, styles). The new_slide_fn is a
# callable that creates a fresh slide with stripe+footer already applied.
# Builders should call new_slide_fn() only if they have data to render;
# otherwise they return without adding a slide.
# =============================================================================


def _slide_title(new_slide, ctx: dict, colors: PptxColors, styles: dict) -> None:
    """Slide 1: Title slide with company name, rating, and analyst info.

    Layout matches the PDF template: brand-colored top panel (58% height)
    with white text, a white rating badge in the top-right with a small
    "RATING" label above the value, and analyst/firm info below.
    """

    slide = new_slide()

    # Brand-colored top panel covering ~58% of the slide
    top_height = int(SLIDE_HEIGHT * 0.58)
    shape = slide.shapes.add_shape(1, Emu(0), Emu(0), SLIDE_WIDTH, top_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors.brand
    shape.line.fill.background()
    # Push behind other shapes
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)

    # "Equity Research" label (top-left, within the blue panel)
    txbox = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(4), Inches(0.4))
    set_text(txbox.text_frame, "EQUITY RESEARCH",
             FontStyle(Pt(11), color=colors.white))

    # Company name (large, bold, white)
    txbox = slide.shapes.add_textbox(Inches(1.0), Inches(1.1), Inches(8.5), Inches(1.2))
    set_text(txbox.text_frame, ctx.get("company_name", ctx.get("ticker", "")),
             FontStyle(Pt(38), bold=True, color=colors.white))

    # Subtitle line: ticker | exchange | sector
    subtitle_parts = [p for p in [ctx.get("ticker"), ctx.get("exchange"), ctx.get("sector")] if p]
    if subtitle_parts:
        txbox = slide.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(8.5), Inches(0.6))
        set_text(txbox.text_frame, " | ".join(subtitle_parts),
                 FontStyle(Pt(15), color=colors.white))

    # Rating badge (top-right: white rounded rect with "RATING" label + value)
    rating = ctx.get("rating", "N/A")
    rating_color = ctx.get("rating_color")
    badge_left = SLIDE_WIDTH - Inches(3.0)
    badge = slide.shapes.add_shape(5, badge_left, Inches(0.7), Inches(2.2), Inches(1.3))
    badge.fill.solid()
    badge.fill.fore_color.rgb = colors.white
    badge.line.fill.background()
    badge.adjustments[0] = 0.08

    tf = badge.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.1)

    # Small "RATING" label
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    para.space_after = Pt(2)
    run = para.add_run()
    run.text = "RATING"
    apply_font(run, FontStyle(Pt(8), color=colors.secondary_text))

    # Large rating value
    para = tf.add_paragraph()
    para.alignment = PP_ALIGN.CENTER
    para.space_before = Pt(2)
    run = para.add_run()
    run.text = rating
    apply_font(run, FontStyle(
        Pt(30), bold=True,
        color=hex_to_pptx(rating_color) if rating_color else colors.brand,
    ))

    # Report date (below the blue panel)
    txbox = slide.shapes.add_textbox(Inches(1.0), Inches(4.7), Inches(5), Inches(0.4))
    set_text(txbox.text_frame, ctx.get("report_date", ""),
             FontStyle(Pt(12), color=colors.dark_text))

    # Analyst info (bottom-left)
    txbox = slide.shapes.add_textbox(Inches(1.0), Inches(5.3), Inches(5), Inches(1.2))
    tf = txbox.text_frame
    set_text(tf, ctx.get("analyst_name", ""), FontStyle(Pt(13), bold=True, color=colors.dark_text))
    if ctx.get("analyst_sector"):
        add_paragraph(tf, ctx["analyst_sector"], FontStyle(Pt(10), color=colors.secondary_text))
    if ctx.get("analyst_email"):
        add_paragraph(tf, ctx["analyst_email"], FontStyle(Pt(10), color=colors.secondary_text))

    # Firm name (bottom-right, brand colored)
    txbox = slide.shapes.add_textbox(
        SLIDE_WIDTH - Inches(4.5), Inches(5.6), Inches(3.8), Inches(0.5),
    )
    set_text(txbox.text_frame, ctx.get("firm_name", ""),
             FontStyle(Pt(16), bold=True, color=colors.brand), PP_ALIGN.RIGHT)


def _slide_toc(new_slide, ctx: dict, colors: PptxColors, styles: dict) -> None:
    """Slide 2: Table of Contents."""
    entries = ctx.get("toc_entries", [])
    if not entries:
        return

    slide = new_slide()
    add_slide_title(slide, "Table of Contents", styles, colors)

    # Two-column layout for entries
    col_width = CONTENT_WIDTH // 2
    per_column = (len(entries) + 1) // 2

    y_start = CONTENT_TOP
    row_height = Inches(0.45)

    for i, entry in enumerate(entries):
        col = 0 if i < per_column else 1
        row = i if col == 0 else i - per_column
        left = CONTENT_LEFT + col * col_width
        top = y_start + row * row_height

        txbox = slide.shapes.add_textbox(left + Inches(0.2), top, col_width - Inches(0.4), row_height)
        set_text(txbox.text_frame, entry["title"], styles["body"])


def _slide_exec_summary(new_slide, ctx: dict, colors: PptxColors, styles: dict) -> None:
    """Slide 3: Executive Summary with headline, metric cards, catalysts, and key risk."""
    exec_data = ctx.get("exec_summary")
    if not exec_data:
        return

    slide = new_slide()
    add_slide_title(slide, "Executive Summary", styles, colors)

    # Headline
    headline = ctx.get("report_headline", "")
    if headline:
        txbox = slide.shapes.add_textbox(CONTENT_LEFT, Inches(1.0), CONTENT_WIDTH, Inches(0.5))
        set_text(txbox.text_frame, headline,
                 FontStyle(Pt(16), bold=True, color=colors.brand))

    # Metric cards row: Rating, Target Price, Current Price, Upside
    card_w = Inches(2.5)
    card_h = Inches(1.1)
    card_y = Inches(1.7)
    card_gap = Inches(0.3)

    cards = [
        ("Rating", exec_data.get("rating", "N/A"), None),
        ("Target Price", _fmt_price(exec_data.get("target_price")), None),
        ("Current Price", _fmt_price(exec_data.get("current_price")), None),
        ("Upside / Downside", _fmt_pct(exec_data.get("upside_pct")),
         _pct_color(exec_data.get("upside_pct"), colors)),
    ]
    for i, (label, value, vc) in enumerate(cards):
        left = CONTENT_LEFT + i * (card_w + card_gap)
        add_metric_card(slide, left, card_y, card_w, card_h, label, value, colors, styles,
                        value_color=vc)

    y = Inches(3.1)

    # Forward snapshot table (forward EPS / P/E for estimated years)
    forward_snapshot = ctx.get("forward_snapshot", [])
    if forward_snapshot:
        snap_headers = ["Year", "Fwd EPS", "Fwd P/E", "EBITDA Margin"]
        n_snap_rows = len(forward_snapshot) + 1
        snap_w = Inches(5.5)
        tbl = create_table(slide, CONTENT_LEFT, y, snap_w, n_snap_rows, len(snap_headers))
        style_table(tbl, colors, styles, header_labels=snap_headers)
        for i, row in enumerate(forward_snapshot):
            ri = i + 1
            set_cell(tbl.cell(ri, 0), str(row.get("year", "")), styles["table_cell"])
            set_cell(tbl.cell(ri, 1), _fmt_number(row.get("eps")), styles["table_cell"], PP_ALIGN.RIGHT)
            set_cell(tbl.cell(ri, 2), _fmt_number(row.get("forward_pe")), styles["table_cell"], PP_ALIGN.RIGHT)
            set_cell(tbl.cell(ri, 3), _fmt_pct_val(row.get("ebitda_margin")), styles["table_cell"], PP_ALIGN.RIGHT)
        y += n_snap_rows * Inches(0.3) + Inches(0.2)

    # Compact bull/base/bear table (right side if snapshot present, else full width)
    scenario = ctx.get("scenario_summary", {})
    if scenario:
        bbb_headers = ["", "Fair Value", "Upside", "Prob"]
        bbb_left = CONTENT_LEFT + Inches(6.0) if forward_snapshot else CONTENT_LEFT
        bbb_w = CONTENT_WIDTH - Inches(6.0) if forward_snapshot else Inches(5.5)
        bbb_top = Inches(3.1) if forward_snapshot else y

        case_rows = []
        for label in ("bull", "base", "bear"):
            case = scenario.get(label, {})
            if case:
                case_rows.append((
                    label.title(),
                    _fmt_price(case.get("fair_value")),
                    _fmt_pct(case.get("upside_pct")),
                    f"{case.get('probability', '')}%",
                ))

        if case_rows:
            tbl = create_table(slide, bbb_left, bbb_top, bbb_w, len(case_rows) + 1, 4)
            style_table(tbl, colors, styles, header_labels=bbb_headers)
            case_colors_map = {"Bull": colors.positive, "Base": colors.brand, "Bear": colors.negative}
            for i, (label, fv, upside, prob) in enumerate(case_rows):
                ri = i + 1
                lbl_color = case_colors_map.get(label, colors.dark_text)
                set_cell(tbl.cell(ri, 0), label, FontStyle(Pt(9), bold=True, color=lbl_color))
                set_cell(tbl.cell(ri, 1), fv, styles["table_cell"], PP_ALIGN.RIGHT)
                set_cell(tbl.cell(ri, 2), upside, styles["table_cell"], PP_ALIGN.RIGHT)
                set_cell(tbl.cell(ri, 3), prob, styles["table_cell"], PP_ALIGN.CENTER)

            if not forward_snapshot:
                y += (len(case_rows) + 1) * Inches(0.3) + Inches(0.2)

    # Catalysts
    catalysts = exec_data.get("catalysts", [])
    if catalysts:
        txbox = slide.shapes.add_textbox(CONTENT_LEFT, y, CONTENT_WIDTH, Inches(0.3))
        set_text(txbox.text_frame, "KEY CATALYSTS", styles["section_label"])
        y += Inches(0.3)
        for cat in catalysts[:3]:
            txbox = slide.shapes.add_textbox(CONTENT_LEFT + Inches(0.2), y, CONTENT_WIDTH - Inches(0.4), Inches(0.3))
            set_text(txbox.text_frame, f"\u2022  {cat}", styles["body_small"])
            y += Inches(0.3)

    # Key risk
    key_risk = exec_data.get("key_risk")
    if key_risk:
        y = max(y + Inches(0.1), Inches(5.5))
        bar = slide.shapes.add_shape(1, CONTENT_LEFT, y, Inches(0.06), Inches(0.55))
        bar.fill.solid()
        bar.fill.fore_color.rgb = colors.negative
        bar.line.fill.background()
        txbox = slide.shapes.add_textbox(CONTENT_LEFT + Inches(0.2), y, CONTENT_WIDTH - Inches(0.3), Inches(0.55))
        tf = txbox.text_frame
        tf.word_wrap = True
        set_text(tf, "KEY RISK", styles["section_label"])
        add_paragraph(tf, key_risk, styles["body_small"])


def _slide_thesis(new_slide, ctx: dict, colors: PptxColors, styles: dict) -> None:
    """Slide 4: Investment Thesis with numbered pillars or fallback bullets."""

    inv_ctx = ctx.get("investment_context", {})
    pillars = inv_ctx.get("thesis_pillars", [])
    market_missing = inv_ctx.get("market_missing", "")
    risk_reward = inv_ctx.get("risk_reward_ratio", "")
    risk_reward_expl = inv_ctx.get("risk_reward_explanation", "")

    # Fall back to LLM-summarized bullets when structured pillars are absent
    if not pillars:
        bullets = ctx.get("thesis_bullets", [])
        if not bullets:
            return
        slide = new_slide()
        add_slide_title(slide, "Investment Thesis", styles, colors)
        add_bullet_list(slide, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, bullets,
                        colors.brand, colors, styles)
        return

    slide = new_slide()
    add_slide_title(slide, "Investment Thesis", styles, colors)

    y = CONTENT_TOP

    # Numbered thesis pillars
    for pillar in pillars[:5]:
        num = pillar.get("number", "")
        statement = pillar.get("statement", "")
        evidence = pillar.get("evidence", "")

        # Pillar number badge (small brand-colored circle)
        badge = slide.shapes.add_shape(
            9, CONTENT_LEFT, y + Inches(0.05), Inches(0.35), Inches(0.35),  # 9 = oval
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = colors.brand
        badge.line.fill.background()
        set_text(badge.text_frame, str(num),
                 FontStyle(Pt(11), bold=True, color=colors.white), PP_ALIGN.CENTER)

        # Statement text
        txbox = slide.shapes.add_textbox(
            CONTENT_LEFT + Inches(0.5), y, CONTENT_WIDTH - Inches(0.5), Inches(0.35),
        )
        txbox.text_frame.word_wrap = True
        set_text(txbox.text_frame, statement,
                 FontStyle(Pt(11), bold=True, color=colors.dark_text))

        # Evidence (smaller, below)
        if evidence:
            y += Inches(0.38)
            txbox = slide.shapes.add_textbox(
                CONTENT_LEFT + Inches(0.5), y, CONTENT_WIDTH - Inches(0.5), Inches(0.3),
            )
            txbox.text_frame.word_wrap = True
            set_text(txbox.text_frame, evidence,
                     FontStyle(Pt(9), color=colors.secondary_text, italic=True))

        y += Inches(0.45)

    # "What the market is missing" paragraph
    if market_missing:
        y += Inches(0.15)
        bar = slide.shapes.add_shape(1, CONTENT_LEFT, y, Inches(0.06), Inches(0.8))
        bar.fill.solid()
        bar.fill.fore_color.rgb = colors.brand
        bar.line.fill.background()

        txbox = slide.shapes.add_textbox(
            CONTENT_LEFT + Inches(0.2), y, CONTENT_WIDTH - Inches(0.3), Inches(0.8),
        )
        tf = txbox.text_frame
        tf.word_wrap = True
        set_text(tf, "WHAT THE MARKET IS MISSING", styles["section_label"])
        add_paragraph(tf, market_missing,
                      FontStyle(Pt(9), color=colors.dark_text))

    # Risk/reward ratio (bottom)
    if risk_reward:
        y_rr = SLIDE_HEIGHT - Inches(1.2)
        card_w = Inches(3.0)
        add_metric_card(slide, CONTENT_LEFT, y_rr, card_w, Inches(0.8),
                        "Risk/Reward", risk_reward, colors, styles)
        if risk_reward_expl:
            txbox = slide.shapes.add_textbox(
                CONTENT_LEFT + card_w + Inches(0.3), y_rr + Inches(0.15),
                CONTENT_WIDTH - card_w - Inches(0.3), Inches(0.6),
            )
            txbox.text_frame.word_wrap = True
            set_text(txbox.text_frame, risk_reward_expl,
                     FontStyle(Pt(9), color=colors.secondary_text))


def _slide_conviction(new_slide, ctx: dict, colors: PptxColors, styles: dict) -> None:
    """Slide 5: Conviction & Rating with overall score and category bars."""
    categories = ctx.get("conviction_categories", [])
    structured = ctx.get("conviction_structured")
    if not categories and not structured:
        return

    slide = new_slide()
    add_slide_title(slide, "Conviction & Rating", styles, colors)

    y = CONTENT_TOP

    # Overall score and confidence metric cards
    if structured:
        overall = structured.get("overall_score", "")
        confidence = structured.get("confidence", "")
        rating = ctx.get("rating", "N/A")

        card_w = Inches(2.8)
        card_h = Inches(1.0)
        gap = Inches(0.3)
        cards_data = [
            ("Overall Score", f"{overall}/100", None),
            ("Confidence", f"{confidence}%" if confidence else "N/A", None),
            ("Rating", rating, hex_to_pptx(ctx["rating_color"]) if ctx.get("rating_color") else None),
        ]
        for i, (label, val, vc) in enumerate(cards_data):
            add_metric_card(slide, CONTENT_LEFT + i * (card_w + gap), y, card_w, card_h,
                            label, val, colors, styles, value_color=vc)
        y += card_h + Inches(0.4)

    # Category score bars
    if categories:
        bar_max_w = Inches(6.0)
        bar_h = Inches(0.35)
        label_w = Inches(2.8)
        score_w = Inches(0.6)
        row_h = Inches(0.55)

        for cat in categories[:7]:
            name = cat.get("name", "")
            score = cat.get("score", 0)
            bar_pct = cat.get("bar_pct", 0)

            # Label
            txbox = slide.shapes.add_textbox(CONTENT_LEFT, y, label_w, bar_h)
            set_text(txbox.text_frame, name, styles["body_small"], PP_ALIGN.RIGHT)

            # Background bar
            bar_left = CONTENT_LEFT + label_w + Inches(0.2)
            bg_bar = slide.shapes.add_shape(1, bar_left, y, bar_max_w, bar_h)
            bg_bar.fill.solid()
            bg_bar.fill.fore_color.rgb = colors.section_bg
            bg_bar.line.fill.background()

            # Foreground bar
            fg_width = max(Inches(0.1), int(bar_max_w * bar_pct / 100))
            fg_bar = slide.shapes.add_shape(1, bar_left, y, fg_width, bar_h)
            fg_bar.fill.solid()
            fg_bar.fill.fore_color.rgb = colors.brand
            fg_bar.line.fill.background()

            # Score value
            txbox = slide.shapes.add_textbox(bar_left + bar_max_w + Inches(0.1), y, score_w, bar_h)
            set_text(txbox.text_frame, f"{score:.1f}", styles["body_small"])

            y += row_h

    # Conviction summary prose
    summary = ctx.get("conviction_summary", "")
    if summary:
        y += Inches(0.2)
        txbox = slide.shapes.add_textbox(CONTENT_LEFT, y, CONTENT_WIDTH, Inches(0.8))
        txbox.text_frame.word_wrap = True
        set_text(txbox.text_frame, summary, styles["disclaimer"])


def _slide_financial_overview(new_slide, ctx: dict, colors: PptxColors, styles: dict) -> None:
    """Slide 6: Financial Overview with key data, statistics, and returns tables."""
    key_data = ctx.get("key_data", [])
    statistics = ctx.get("statistics", [])
    if not key_data and not statistics:
        return

    slide = new_slide()
    add_slide_title(slide, "Financial Overview", styles, colors)

    half_w = (CONTENT_WIDTH - Inches(0.4)) // 2

    # Key Data table (left)
    if key_data:
        rows = len(key_data) + 1
        tbl = create_table(slide, CONTENT_LEFT, CONTENT_TOP, half_w, rows, 2)
        style_table(tbl, colors, styles, header_labels=["Metric", "Value"])
        for i, item in enumerate(key_data):
            set_cell(tbl.cell(i + 1, 0), item.get("label", ""), styles["table_cell"])
            set_cell(tbl.cell(i + 1, 1), str(item.get("value", "")), styles["table_cell"], PP_ALIGN.RIGHT)

    # Statistics table (right)
    if statistics:
        rows = len(statistics) + 1
        right_left = CONTENT_LEFT + half_w + Inches(0.4)
        tbl = create_table(slide, right_left, CONTENT_TOP, half_w, rows, 2)
        style_table(tbl, colors, styles, header_labels=["Statistic", "Value"])
        for i, item in enumerate(statistics):
            set_cell(tbl.cell(i + 1, 0), item.get("label", ""), styles["table_cell"])
            set_cell(tbl.cell(i + 1, 1), str(item.get("value", "")), styles["table_cell"], PP_ALIGN.RIGHT)

    # Returns table (below both)
    returns = ctx.get("returns", [])
    if returns:
        returns_top = CONTENT_TOP + Inches(3.2)
        cols = len(returns) + 1
        tbl = create_table(slide, CONTENT_LEFT, returns_top, CONTENT_WIDTH, 2, cols)
        style_table(tbl, colors, styles, header_labels=["Period Returns"] + [r.get("period", "") for r in returns])
        set_cell(tbl.cell(1, 0), "Return", styles["table_cell"])
        for i, ret in enumerate(returns):
            display = ret.get("display", "N/A")
            val = ret.get("value")
            vc = _pct_color(val, colors) if val is not None else None
            cell_style = FontStyle(Pt(9), bold=True, color=vc) if vc else styles["table_cell"]
            set_cell(tbl.cell(1, i + 1), display, cell_style, PP_ALIGN.CENTER)


def _slide_football_field(new_slide, ctx: dict, colors: PptxColors, styles: dict) -> None:
    """Slide 7: Valuation Football Field as native PowerPoint shapes.

    Renders horizontal range bars, mid-point diamonds, and a dashed
    current-price line using editable shapes instead of a raster PNG.
    Falls back to a table if fewer than one range is provided.
    """
    ranges = ctx.get("football_ranges", [])
    if not ranges:
        return

    current_price = ctx.get("football_current_price")

    slide = new_slide()
    add_slide_title(slide, "Valuation Football Field", styles, colors)

    _render_football_field_native(slide, ranges, current_price, colors, styles)


def _slide_football_field_table(new_slide, ctx: dict, colors: PptxColors, styles: dict) -> None:
    """Fallback: render football field as a text table when matplotlib fails."""
    ranges = ctx.get("football_ranges", [])
    if not ranges:
        return

    slide = new_slide()
    add_slide_title(slide, "Valuation Football Field", styles, colors)

    rows = len(ranges) + 1
    tbl = create_table(slide, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, rows, 4)
    style_table(tbl, colors, styles, header_labels=["Method", "Low", "Mid", "High"])
    for i, r in enumerate(ranges):
        set_cell(tbl.cell(i + 1, 0), r.get("method", ""), styles["table_cell"])
        set_cell(tbl.cell(i + 1, 1), _fmt_price(r.get("low")), styles["table_cell"], PP_ALIGN.RIGHT)
        set_cell(tbl.cell(i + 1, 2), _fmt_price(r.get("mid")), styles["table_cell"], PP_ALIGN.RIGHT)
        set_cell(tbl.cell(i + 1, 3), _fmt_price(r.get("high")), styles["table_cell"], PP_ALIGN.RIGHT)


def _slide_dcf(new_slide, ctx: dict, colors: PptxColors, styles: dict) -> None:
    """Slide 8: DCF Valuation with metric cards."""
    dcf = ctx.get("dcf_summary")
    if not dcf:
        return

    slide = new_slide()
    add_slide_title(slide, "DCF Valuation", styles, colors)

    card_w = Inches(2.8)
    card_h = Inches(1.1)
    gap = Inches(0.3)
    cards_per_row = 4
    y = CONTENT_TOP

    card_data = [
        ("Fair Value / Share", _fmt_price(dcf.get("fair_value")),
         _pct_color(dcf.get("upside_pct"), colors)),
        ("Upside / Downside", _fmt_pct(dcf.get("upside_pct")),
         _pct_color(dcf.get("upside_pct"), colors)),
        ("WACC", _fmt_pct_val(dcf.get("wacc")), None),
        ("Terminal Growth", _fmt_pct_val(dcf.get("terminal_growth_rate")), None),
        ("Enterprise Value", _fmt_large_number(dcf.get("enterprise_value")), None),
        ("TV % of EV", _fmt_pct_val(dcf.get("tv_pct_of_ev")), None),
    ]

    for i, (label, value, vc) in enumerate(card_data):
        col = i % cards_per_row
        row = i // cards_per_row
        left = CONTENT_LEFT + col * (card_w + gap)
        top = y + row * (card_h + gap)
        add_metric_card(slide, left, top, card_w, card_h, label, value, colors, styles,
                        value_color=vc)


def _slide_scenarios(new_slide, ctx: dict, colors: PptxColors, styles: dict) -> None:
    """Slide 9: Scenario Analysis with bull/base/bear columns."""
    scenario = ctx.get("scenario_summary")
    if not scenario:
        return

    slide = new_slide()
    add_slide_title(slide, "Scenario Analysis", styles, colors)

    col_w = (CONTENT_WIDTH - Inches(0.6)) // 3
    col_gap = Inches(0.3)

    # Map each case to its gradient background color and accent
    case_configs = [
        ("bull", "Bull Case", colors.positive, colors.bull_bg),
        ("base", "Base Case", colors.brand, colors.section_bg),
        ("bear", "Bear Case", colors.negative, colors.bear_bg),
    ]

    col_h = CONTENT_HEIGHT - Inches(0.1)

    for i, (key, title, color, bg_color) in enumerate(case_configs):
        case_data = scenario.get(key, {})
        if not case_data:
            continue

        left = CONTENT_LEFT + i * (col_w + col_gap)
        y = CONTENT_TOP

        # Subtle gradient background panel (colored to white, top to bottom)
        panel = slide.shapes.add_shape(5, left, y, col_w, col_h)
        apply_gradient_fill(panel, bg_color, hex_to_pptx("FFFFFF"), angle_degrees=90.0)
        panel.line.fill.background()
        panel.adjustments[0] = 0.03
        slide.shapes._spTree.remove(panel._element)
        slide.shapes._spTree.insert(2, panel._element)

        # Colored header bar
        header = slide.shapes.add_shape(1, left, y, col_w, Inches(0.4))
        header.fill.solid()
        header.fill.fore_color.rgb = color
        header.line.fill.background()
        set_text(header.text_frame, title,
                 FontStyle(Pt(12), bold=True, color=colors.white), PP_ALIGN.CENTER)
        y += Inches(0.5)

        # Fair value card
        add_metric_card(slide, left, y, col_w, Inches(1.0),
                        "Fair Value", _fmt_price(case_data.get("fair_value")),
                        colors, styles)
        y += Inches(1.15)

        # Probability and upside
        txbox = slide.shapes.add_textbox(left, y, col_w, Inches(0.3))
        prob = case_data.get("probability", "")
        prob_text = f"Probability: {prob}%" if prob else ""
        set_text(txbox.text_frame, prob_text, styles["body_small"], PP_ALIGN.CENTER)
        y += Inches(0.35)

        upside = case_data.get("upside_pct")
        if upside is not None:
            txbox = slide.shapes.add_textbox(left, y, col_w, Inches(0.3))
            set_text(txbox.text_frame, f"Upside: {_fmt_pct(upside)}",
                     FontStyle(Pt(10), color=_pct_color(upside, colors)), PP_ALIGN.CENTER)
            y += Inches(0.4)

        # Assumptions
        assumptions = case_data.get("assumptions", {})
        if assumptions:
            txbox = slide.shapes.add_textbox(left + Inches(0.1), y, col_w - Inches(0.2), Inches(1.5))
            tf = txbox.text_frame
            tf.word_wrap = True
            set_text(tf, "Assumptions:", FontStyle(Pt(9), bold=True, color=colors.label_text))
            for a_key, a_val in assumptions.items():
                label = a_key.replace("_", " ").title()
                add_paragraph(tf, f"{label}: {_fmt_assumption(a_val)}",
                              FontStyle(Pt(8), color=colors.secondary_text, italic=True))

    # Bottom row: weighted fair value + risk/reward ratio
    weighted = scenario.get("weighted_fair_value")
    inv_ctx = ctx.get("investment_context", {})
    risk_reward = inv_ctx.get("risk_reward_ratio", "")

    y = SLIDE_HEIGHT - Inches(1.5)
    if weighted:
        wfv_w = Inches(3.0) if risk_reward else Inches(3.6)
        add_metric_card(slide, CONTENT_LEFT + Inches(0.5), y,
                        wfv_w, Inches(0.9),
                        "Weighted Fair Value", _fmt_price(weighted), colors, styles)

    if risk_reward:
        rr_left = CONTENT_LEFT + Inches(4.0) if weighted else CONTENT_LEFT + CONTENT_WIDTH // 2 - Inches(1.5)
        add_metric_card(slide, rr_left, y,
                        Inches(3.0), Inches(0.9),
                        "Risk / Reward", risk_reward, colors, styles)


def _slide_sensitivity(new_slide, ctx: dict, colors: PptxColors, styles: dict) -> None:
    """Slide 10: Sensitivity Analysis grids with colored cells."""
    grid = ctx.get("sensitivity_grid")
    growth_grid = ctx.get("growth_margin_grid")
    if not grid and not growth_grid:
        return

    slide = new_slide()
    add_slide_title(slide, "Sensitivity Analysis", styles, colors)

    y = CONTENT_TOP

    if grid:
        y = _render_sensitivity_grid(
            slide, grid, "WACC vs Terminal Growth Rate",
            CONTENT_LEFT, y, CONTENT_WIDTH, colors, styles,
        )
        y += Inches(0.4)

    if growth_grid:
        _render_sensitivity_grid(
            slide, growth_grid, "Revenue Growth vs Operating Margin",
            CONTENT_LEFT, y, CONTENT_WIDTH, colors, styles,
        )


def _render_sensitivity_grid(
    slide, grid: dict, title: str,
    left: int, top: int, width: int,
    colors: PptxColors, styles: dict,
) -> int:
    """Render a single sensitivity grid table. Returns the bottom edge position."""
    row_values = grid.get("row_values", [])
    col_values = grid.get("col_values", [])
    data_grid = grid.get("grid", [])
    cell_classes = grid.get("cell_classes", [])
    base_row = grid.get("base_row_idx")
    base_col = grid.get("base_col_idx")

    if not row_values or not col_values:
        return top

    # Section label
    txbox = slide.shapes.add_textbox(left, top, width, Inches(0.3))
    set_text(txbox.text_frame, title, styles["section_label"])
    top += Inches(0.35)

    rows = len(row_values) + 1  # +1 for header
    cols = len(col_values) + 1  # +1 for row labels
    tbl = create_table(slide, left, top, width, rows, cols)

    # Header row
    set_cell(tbl.cell(0, 0), "", styles["table_header"])
    tbl.cell(0, 0).fill.solid()
    tbl.cell(0, 0).fill.fore_color.rgb = colors.brand
    for ci, cv in enumerate(col_values):
        cell = tbl.cell(0, ci + 1)
        cell.fill.solid()
        cell.fill.fore_color.rgb = colors.brand
        set_cell(cell, _fmt_pct_val(cv), styles["table_header"], PP_ALIGN.CENTER)

    # Data rows
    for ri, rv in enumerate(row_values):
        # Row label
        label_cell = tbl.cell(ri + 1, 0)
        label_cell.fill.solid()
        label_cell.fill.fore_color.rgb = colors.section_bg
        set_cell(label_cell, _fmt_pct_val(rv),
                 FontStyle(Pt(9), bold=True, color=colors.dark_text))

        for ci in range(len(col_values)):
            cell = tbl.cell(ri + 1, ci + 1)
            val = data_grid[ri][ci] if ri < len(data_grid) and ci < len(data_grid[ri]) else ""

            # Color based on cell class
            cls = ""
            if cell_classes and ri < len(cell_classes) and ci < len(cell_classes[ri]):
                cls = cell_classes[ri][ci]

            if cls == "base" or (base_row is not None and ri == base_row and base_col is not None and ci == base_col):
                cell.fill.solid()
                cell.fill.fore_color.rgb = colors.base_case_bg
                set_cell(cell, _fmt_price(val),
                         FontStyle(Pt(9), bold=True, color=colors.dark_text), PP_ALIGN.CENTER)
            elif cls == "positive":
                cell.fill.solid()
                cell.fill.fore_color.rgb = _lighten(colors.positive, 0.85)
                set_cell(cell, _fmt_price(val), styles["table_cell"], PP_ALIGN.CENTER)
            elif cls == "negative":
                cell.fill.solid()
                cell.fill.fore_color.rgb = _lighten(colors.negative, 0.85)
                set_cell(cell, _fmt_price(val), styles["table_cell"], PP_ALIGN.CENTER)
            else:
                set_cell(cell, _fmt_price(val), styles["table_cell"], PP_ALIGN.CENTER)

    row_height = Inches(0.35)
    return top + rows * row_height


def _slide_earnings(new_slide, ctx: dict, colors: PptxColors, styles: dict) -> None:
    """Slide 11: Earnings Model table."""
    earnings = ctx.get("earnings_table")
    if not earnings or not earnings.get("rows"):
        return

    slide = new_slide()
    add_slide_title(slide, "Earnings Model", styles, colors)

    rows_data = earnings["rows"]
    # Columns: FY, Revenue, Rev Growth, EBITDA, EBITDA Margin, EPS
    headers = ["Fiscal Year", "Revenue", "Rev Growth", "EBITDA", "EBITDA Margin", "EPS"]
    n_rows = len(rows_data) + 1
    n_cols = len(headers)

    tbl = create_table(slide, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, n_rows, n_cols)
    style_table(tbl, colors, styles, header_labels=headers)

    for i, row in enumerate(rows_data):
        ri = i + 1
        is_est = row.get("is_estimate", False)

        set_cell(tbl.cell(ri, 0), str(row.get("fiscal_year", "")), styles["table_cell"])
        set_cell(tbl.cell(ri, 1), _fmt_large_number(row.get("revenue")), styles["table_cell"], PP_ALIGN.RIGHT)

        # Revenue growth with color
        rg = row.get("revenue_growth")
        rg_color = _pct_color(rg, colors)
        rg_style = FontStyle(Pt(9), color=rg_color) if rg_color else styles["table_cell"]
        set_cell(tbl.cell(ri, 2), _fmt_pct_val(rg), rg_style, PP_ALIGN.RIGHT)

        set_cell(tbl.cell(ri, 3), _fmt_large_number(row.get("ebitda")), styles["table_cell"], PP_ALIGN.RIGHT)
        set_cell(tbl.cell(ri, 4), _fmt_pct_val(row.get("ebitda_margin")), styles["table_cell"], PP_ALIGN.RIGHT)
        set_cell(tbl.cell(ri, 5), _fmt_number(row.get("eps")), styles["table_cell"], PP_ALIGN.RIGHT)

        # Shade estimate rows
        if is_est:
            for ci in range(n_cols):
                cell = tbl.cell(ri, ci)
                cell.fill.solid()
                cell.fill.fore_color.rgb = colors.estimate_row_bg

    # Earnings takeaway
    takeaway = ctx.get("earnings_takeaway")
    if takeaway:
        y = CONTENT_TOP + (n_rows + 1) * Inches(0.35)
        # Yellow callout box
        bar = slide.shapes.add_shape(1, CONTENT_LEFT, y, Inches(0.06), Inches(0.5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = colors.brand
        bar.line.fill.background()
        txbox = slide.shapes.add_textbox(CONTENT_LEFT + Inches(0.2), y, CONTENT_WIDTH - Inches(0.3), Inches(0.5))
        txbox.text_frame.word_wrap = True
        set_text(txbox.text_frame, takeaway, styles["body_small"])


def _slide_comps(new_slide, ctx: dict, colors: PptxColors, styles: dict) -> None:
    """Slide 12: Comparable Companies with implied values and peer table."""
    comp = ctx.get("comp_summary")
    if not comp:
        return

    slide = new_slide()
    add_slide_title(slide, "Comparable Companies", styles, colors)

    y = CONTENT_TOP

    # Implied value metric cards
    implied = comp.get("implied_values", {})
    if implied:
        card_w = Inches(3.0)
        card_h = Inches(0.9)
        gap = Inches(0.3)
        iv_cards = [
            ("P/E Implied", _fmt_price(implied.get("pe_implied"))),
            ("EV/EBITDA Implied", _fmt_price(implied.get("ev_ebitda_implied"))),
            ("EV/Revenue Implied", _fmt_price(implied.get("ev_revenue_implied"))),
        ]
        for i, (label, val) in enumerate(iv_cards):
            add_metric_card(slide, CONTENT_LEFT + i * (card_w + gap), y, card_w, card_h,
                            label, val, colors, styles)
        y += card_h + Inches(0.3)

    # Peer table
    peers = comp.get("peers", [])
    target = comp.get("target", {})
    medians = comp.get("medians", {})

    all_rows = []
    if target:
        all_rows.append(target)
    all_rows.extend(peers)

    if not all_rows:
        return

    headers = ["Company", "Mkt Cap", "P/E", "Fwd P/E", "EV/EBITDA", "EV/Rev", "Gross Margin", "Rev Growth"]
    n_rows = len(all_rows) + 2  # header + data + median
    tbl = create_table(slide, CONTENT_LEFT, y, CONTENT_WIDTH, n_rows, len(headers))
    style_table(tbl, colors, styles, header_labels=headers)

    for i, row in enumerate(all_rows):
        ri = i + 1
        is_target = (target and i == 0)

        set_cell(tbl.cell(ri, 0), row.get("ticker", row.get("name", "")), styles["table_cell"])
        set_cell(tbl.cell(ri, 1), _fmt_large_number(row.get("market_cap")), styles["table_cell"], PP_ALIGN.RIGHT)
        set_cell(tbl.cell(ri, 2), _fmt_number(row.get("pe_ratio")), styles["table_cell"], PP_ALIGN.RIGHT)
        set_cell(tbl.cell(ri, 3), _fmt_number(row.get("forward_pe")), styles["table_cell"], PP_ALIGN.RIGHT)
        set_cell(tbl.cell(ri, 4), _fmt_number(row.get("ev_to_ebitda")), styles["table_cell"], PP_ALIGN.RIGHT)
        set_cell(tbl.cell(ri, 5), _fmt_number(row.get("ev_to_revenue")), styles["table_cell"], PP_ALIGN.RIGHT)
        set_cell(tbl.cell(ri, 6), _fmt_pct_val(row.get("gross_margin")), styles["table_cell"], PP_ALIGN.RIGHT)
        set_cell(tbl.cell(ri, 7), _fmt_pct_val(row.get("revenue_growth")), styles["table_cell"], PP_ALIGN.RIGHT)

        # Highlight target row
        if is_target:
            for ci in range(len(headers)):
                tbl.cell(ri, ci).fill.solid()
                tbl.cell(ri, ci).fill.fore_color.rgb = colors.base_case_bg

    # Median row
    if medians:
        mi = n_rows - 1
        set_cell(tbl.cell(mi, 0), "Median",
                 FontStyle(Pt(9), bold=True, italic=True, color=colors.dark_text))
        set_cell(tbl.cell(mi, 1), "", styles["table_cell"])
        set_cell(tbl.cell(mi, 2), _fmt_number(medians.get("pe_ratio")), styles["table_cell"], PP_ALIGN.RIGHT)
        set_cell(tbl.cell(mi, 3), _fmt_number(medians.get("forward_pe")), styles["table_cell"], PP_ALIGN.RIGHT)
        set_cell(tbl.cell(mi, 4), _fmt_number(medians.get("ev_to_ebitda")), styles["table_cell"], PP_ALIGN.RIGHT)
        set_cell(tbl.cell(mi, 5), _fmt_number(medians.get("ev_to_revenue")), styles["table_cell"], PP_ALIGN.RIGHT)
        set_cell(tbl.cell(mi, 6), _fmt_pct_val(medians.get("gross_margin")), styles["table_cell"], PP_ALIGN.RIGHT)
        set_cell(tbl.cell(mi, 7), _fmt_pct_val(medians.get("revenue_growth")), styles["table_cell"], PP_ALIGN.RIGHT)


def _slide_precedents(new_slide, ctx: dict, colors: PptxColors, styles: dict) -> None:
    """Slide 13: Precedent M&A Transactions."""
    prec = ctx.get("precedent_summary")
    if not prec:
        return

    deals = prec.get("deals", [])
    if not deals:
        return

    slide = new_slide()
    add_slide_title(slide, "Precedent M&A Transactions", styles, colors)

    y = CONTENT_TOP

    # Implied value cards
    implied = prec.get("implied_values", {})
    if implied:
        card_w = Inches(3.0)
        card_h = Inches(0.9)
        gap = Inches(0.3)
        cards = [
            ("EV/Revenue Implied", _fmt_price(implied.get("ev_revenue"))),
            ("EV/EBITDA Implied", _fmt_price(implied.get("ev_ebitda"))),
        ]
        for i, (label, val) in enumerate(cards):
            add_metric_card(slide, CONTENT_LEFT + i * (card_w + gap), y, card_w, card_h,
                            label, val, colors, styles)
        y += card_h + Inches(0.3)

    # Deals table
    headers = ["Date", "Acquirer", "Target", "EV/Rev", "EV/EBITDA", "Premium"]
    n_rows = len(deals) + 2  # header + data + median
    tbl = create_table(slide, CONTENT_LEFT, y, CONTENT_WIDTH, n_rows, len(headers))
    style_table(tbl, colors, styles, header_labels=headers)

    for i, deal in enumerate(deals):
        ri = i + 1
        set_cell(tbl.cell(ri, 0), str(deal.get("date", "")), styles["table_cell"])
        set_cell(tbl.cell(ri, 1), str(deal.get("acquirer", "")), styles["table_cell"])
        set_cell(tbl.cell(ri, 2), str(deal.get("target", "")), styles["table_cell"])
        set_cell(tbl.cell(ri, 3), _fmt_number(deal.get("ev_to_revenue")), styles["table_cell"], PP_ALIGN.RIGHT)
        set_cell(tbl.cell(ri, 4), _fmt_number(deal.get("ev_to_ebitda")), styles["table_cell"], PP_ALIGN.RIGHT)
        set_cell(tbl.cell(ri, 5), _fmt_pct_val(deal.get("premium_pct")), styles["table_cell"], PP_ALIGN.RIGHT)

    # Median row
    medians = prec.get("medians", {})
    if medians:
        mi = n_rows - 1
        set_cell(tbl.cell(mi, 0), "Median",
                 FontStyle(Pt(9), bold=True, italic=True, color=colors.dark_text))
        set_cell(tbl.cell(mi, 1), "", styles["table_cell"])
        set_cell(tbl.cell(mi, 2), "", styles["table_cell"])
        set_cell(tbl.cell(mi, 3), _fmt_number(medians.get("ev_revenue")), styles["table_cell"], PP_ALIGN.RIGHT)
        set_cell(tbl.cell(mi, 4), _fmt_number(medians.get("ev_ebitda")), styles["table_cell"], PP_ALIGN.RIGHT)
        set_cell(tbl.cell(mi, 5), _fmt_pct_val(medians.get("premium")), styles["table_cell"], PP_ALIGN.RIGHT)

    # Takeaway
    takeaway = ctx.get("precedent_takeaway")
    if takeaway:
        y_after = y + n_rows * Inches(0.35) + Inches(0.2)
        bar = slide.shapes.add_shape(1, CONTENT_LEFT, y_after, Inches(0.06), Inches(0.5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = colors.brand
        bar.line.fill.background()
        txbox = slide.shapes.add_textbox(CONTENT_LEFT + Inches(0.2), y_after,
                                         CONTENT_WIDTH - Inches(0.3), Inches(0.5))
        txbox.text_frame.word_wrap = True
        set_text(txbox.text_frame, takeaway, styles["body_small"])


def _slide_catalysts_risks(new_slide, ctx: dict, colors: PptxColors, styles: dict) -> None:
    """Catalysts timeline and quantified risks from structured investment context."""

    inv_ctx = ctx.get("investment_context", {})
    catalysts = inv_ctx.get("catalysts", [])
    risks = inv_ctx.get("quantified_risks", [])

    if not catalysts and not risks:
        return

    slide = new_slide()
    add_slide_title(slide, "Catalysts & Risk Quantification", styles, colors)

    y = CONTENT_TOP

    # Catalyst timeline table
    if catalysts:
        txbox = slide.shapes.add_textbox(CONTENT_LEFT, y, CONTENT_WIDTH, Inches(0.3))
        set_text(txbox.text_frame, "CATALYST TIMELINE", styles["section_label"])
        y += Inches(0.3)

        cat_headers = ["Event", "Expected Date", "Expected Impact"]
        n_rows = min(len(catalysts), 6) + 1
        tbl = create_table(slide, CONTENT_LEFT, y, CONTENT_WIDTH, n_rows, len(cat_headers))
        style_table(tbl, colors, styles, header_labels=cat_headers)

        for i, cat in enumerate(catalysts[:6]):
            ri = i + 1
            set_cell(tbl.cell(ri, 0), cat.get("event", ""), styles["table_cell"])
            set_cell(tbl.cell(ri, 1), cat.get("expected_date", ""), styles["table_cell"], PP_ALIGN.CENTER)
            set_cell(tbl.cell(ri, 2), cat.get("expected_impact", ""), styles["table_cell"])

        y += n_rows * Inches(0.35) + Inches(0.3)

    # Quantified risks table
    if risks:
        txbox = slide.shapes.add_textbox(CONTENT_LEFT, y, CONTENT_WIDTH, Inches(0.3))
        set_text(txbox.text_frame, "QUANTIFIED RISK FACTORS", styles["section_label"])
        y += Inches(0.3)

        risk_headers = ["Risk", "Key Metric", "Sensitivity"]
        n_rows = min(len(risks), 5) + 1
        tbl = create_table(slide, CONTENT_LEFT, y, CONTENT_WIDTH, n_rows, len(risk_headers))
        style_table(tbl, colors, styles, header_labels=risk_headers)

        for i, risk in enumerate(risks[:5]):
            ri = i + 1
            set_cell(tbl.cell(ri, 0), risk.get("risk", ""), styles["table_cell"])
            set_cell(tbl.cell(ri, 1), risk.get("current_value", ""), styles["table_cell"], PP_ALIGN.CENTER)
            set_cell(tbl.cell(ri, 2), risk.get("sensitivity", ""), styles["table_cell"])


def _slide_monitoring(new_slide, ctx: dict, colors: PptxColors, styles: dict) -> None:
    """Monitoring KPIs with bull/bear thresholds from structured investment context."""
    inv_ctx = ctx.get("investment_context", {})
    kpis = inv_ctx.get("monitoring_kpis", [])

    if not kpis:
        return

    slide = new_slide()
    add_slide_title(slide, "Thesis Monitoring Framework", styles, colors)

    y = CONTENT_TOP

    headers = ["KPI", "Current", "Bull Threshold", "Bear Threshold", "Frequency"]
    n_rows = min(len(kpis), 6) + 1
    tbl = create_table(slide, CONTENT_LEFT, y, CONTENT_WIDTH, n_rows, len(headers))
    style_table(tbl, colors, styles, header_labels=headers)

    for i, kpi in enumerate(kpis[:6]):
        ri = i + 1
        set_cell(tbl.cell(ri, 0), kpi.get("kpi", ""), styles["table_cell"])
        set_cell(tbl.cell(ri, 1), kpi.get("current_value", ""), styles["table_cell"], PP_ALIGN.CENTER)

        # Bull threshold in green
        bull_cell = tbl.cell(ri, 2)
        set_cell(bull_cell, kpi.get("bull_threshold", ""),
                 FontStyle(Pt(9), color=colors.positive), PP_ALIGN.CENTER)

        # Bear threshold in red
        bear_cell = tbl.cell(ri, 3)
        set_cell(bear_cell, kpi.get("bear_threshold", ""),
                 FontStyle(Pt(9), color=colors.negative), PP_ALIGN.CENTER)

        set_cell(tbl.cell(ri, 4), kpi.get("frequency", ""), styles["table_cell"], PP_ALIGN.CENTER)


def _slide_strategy(new_slide, ctx: dict, colors: PptxColors, styles: dict) -> None:
    """Slide 14: Strategic Assessment (2x2 SWOT grid or fallback bullets)."""
    grid = ctx.get("strategy_grid")
    bullets = ctx.get("strategy_bullets", [])

    if not grid and not bullets:
        return

    slide = new_slide()
    add_slide_title(slide, "Strategic Assessment", styles, colors)

    if grid:
        # 2x2 grid of colored boxes
        half_w = (CONTENT_WIDTH - Inches(0.3)) // 2
        half_h = (CONTENT_HEIGHT - Inches(0.3)) // 2

        quadrants = [
            ("advantages", "Advantages", colors.swot_strengths_border, colors.swot_strengths_bg, 0, 0),
            ("vulnerabilities", "Vulnerabilities", colors.swot_weaknesses_border, colors.swot_weaknesses_bg, 1, 0),
            ("tailwinds", "Tailwinds", colors.swot_opportunities_border, colors.swot_opportunities_bg, 0, 1),
            ("headwinds", "Headwinds", colors.swot_threats_border, colors.swot_threats_bg, 1, 1),
        ]

        for key, title, border_color, bg_color, col, row in quadrants:
            content = grid.get(key, "")
            if not content:
                continue

            left = CONTENT_LEFT + col * (half_w + Inches(0.3))
            top = CONTENT_TOP + row * (half_h + Inches(0.3))

            # Background box
            box = slide.shapes.add_shape(1, left, top, half_w, half_h)
            box.fill.solid()
            box.fill.fore_color.rgb = bg_color
            box.line.color.rgb = border_color
            box.line.width = Pt(1.5)

            # Left accent bar
            bar = slide.shapes.add_shape(1, left, top, Inches(0.06), half_h)
            bar.fill.solid()
            bar.fill.fore_color.rgb = border_color
            bar.line.fill.background()

            # Title and content
            txbox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.08),
                                             half_w - Inches(0.3), half_h - Inches(0.15))
            tf = txbox.text_frame
            tf.word_wrap = True
            set_text(tf, title, FontStyle(Pt(11), bold=True, color=border_color))

            # Parse HTML list items into plain text bullets
            items = _parse_html_list(content)
            for item in items[:5]:
                add_paragraph(tf, f"\u2022  {item}", styles["body_small"])
    else:
        add_bullet_list(slide, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, bullets,
                        colors.brand, colors, styles)


def _slide_industry(new_slide, ctx: dict, colors: PptxColors, styles: dict) -> None:
    """Slide 15: Industry Dynamics bullets."""
    bullets = ctx.get("external_bullets", [])
    if not bullets:
        return
    slide = new_slide()
    add_slide_title(slide, "Industry Dynamics", styles, colors)
    add_bullet_list(slide, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, bullets,
                    colors.brand, colors, styles)


def _slide_moat(new_slide, ctx: dict, colors: PptxColors, styles: dict) -> None:
    """Slide 16: Competitive Moat bullets."""
    bullets = ctx.get("moat_bullets", [])
    if not bullets:
        return
    slide = new_slide()
    add_slide_title(slide, "Competitive Moat", styles, colors)
    add_bullet_list(slide, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, bullets,
                    colors.brand, colors, styles)


def _slide_investment_cases(new_slide, ctx: dict, colors: PptxColors, styles: dict) -> None:
    """Slide 17: Bull and Bear cases side by side with gradient backgrounds.

    Each case gets a gradient box (green-to-white for bull, red-to-white for
    bear) with a thick left accent bar, matching the PDF template's case-box
    design.
    """
    bull = ctx.get("bull_bullets", [])
    bear = ctx.get("bear_bullets", [])
    if not bull and not bear:
        return

    slide = new_slide()
    add_slide_title(slide, "Investment Cases", styles, colors)

    half_w = (CONTENT_WIDTH - Inches(0.5)) // 2
    box_h = CONTENT_HEIGHT - Inches(0.1)
    y = CONTENT_TOP

    # Bull case
    if bull:
        # Gradient background box (light green to white, left to right)
        bg = slide.shapes.add_shape(5, CONTENT_LEFT, y, half_w, box_h)
        apply_gradient_fill(bg, colors.bull_bg, hex_to_pptx("FFFFFF"), angle_degrees=0.0)
        bg.line.fill.background()
        bg.adjustments[0] = 0.03
        # Push behind text
        slide.shapes._spTree.remove(bg._element)
        slide.shapes._spTree.insert(2, bg._element)

        # Left accent bar
        bar = slide.shapes.add_shape(1, CONTENT_LEFT, y, Inches(0.08), box_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = colors.bull_border
        bar.line.fill.background()

        # Header
        txbox = slide.shapes.add_textbox(CONTENT_LEFT + Inches(0.2), y + Inches(0.1),
                                         half_w - Inches(0.3), Inches(0.35))
        set_text(txbox.text_frame, "Bull Case",
                 FontStyle(Pt(14), bold=True, color=colors.positive))

        add_bullet_list(slide, CONTENT_LEFT + Inches(0.2), y + Inches(0.5),
                        half_w - Inches(0.3), bull,
                        colors.bull_border, colors, styles)

    # Bear case
    if bear:
        right_left = CONTENT_LEFT + half_w + Inches(0.5)

        # Gradient background box (light red to white, left to right)
        bg = slide.shapes.add_shape(5, right_left, y, half_w, box_h)
        apply_gradient_fill(bg, colors.bear_bg, hex_to_pptx("FFFFFF"), angle_degrees=0.0)
        bg.line.fill.background()
        bg.adjustments[0] = 0.03
        slide.shapes._spTree.remove(bg._element)
        slide.shapes._spTree.insert(2, bg._element)

        # Left accent bar
        bar = slide.shapes.add_shape(1, right_left, y, Inches(0.08), box_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = colors.bear_border
        bar.line.fill.background()

        # Header
        txbox = slide.shapes.add_textbox(right_left + Inches(0.2), y + Inches(0.1),
                                         half_w - Inches(0.3), Inches(0.35))
        set_text(txbox.text_frame, "Bear Case",
                 FontStyle(Pt(14), bold=True, color=colors.negative))

        add_bullet_list(slide, right_left + Inches(0.2), y + Inches(0.5),
                        half_w - Inches(0.3), bear,
                        colors.bear_border, colors, styles)


def _slide_sources(new_slide, ctx: dict, colors: PptxColors, styles: dict) -> None:
    """Slide 18: Sources & References."""
    research = ctx.get("research_sources", [])
    analyst = ctx.get("analyst_sources", [])
    if not research and not analyst:
        return

    slide = new_slide()
    add_slide_title(slide, "Sources & References", styles, colors)

    y = CONTENT_TOP

    if research:
        txbox = slide.shapes.add_textbox(CONTENT_LEFT, y, CONTENT_WIDTH, Inches(0.3))
        set_text(txbox.text_frame, "Research Sources", styles["section_label"])
        y += Inches(0.35)

        # Two-column layout for sources
        half_w = (CONTENT_WIDTH - Inches(0.3)) // 2
        per_col = (len(research) + 1) // 2

        for i, src in enumerate(research):
            col = 0 if i < per_col else 1
            row = i if col == 0 else i - per_col
            left = CONTENT_LEFT + col * (half_w + Inches(0.3))
            src_y = y + row * Inches(0.3)

            title = src.get("title", "")
            date = src.get("date", "")
            date_str = f" ({date})" if date else ""
            text = f"[{src.get('id', i + 1)}] {title}{date_str}"

            txbox = slide.shapes.add_textbox(left, src_y, half_w, Inches(0.3))
            set_text(txbox.text_frame, text,
                     FontStyle(Pt(8), color=colors.secondary_text))

        y += (per_col + 1) * Inches(0.3)

    if analyst:
        txbox = slide.shapes.add_textbox(CONTENT_LEFT, y, CONTENT_WIDTH, Inches(0.3))
        set_text(txbox.text_frame, "Analyst Notes", styles["section_label"])
        y += Inches(0.35)

        for src in analyst:
            txbox = slide.shapes.add_textbox(CONTENT_LEFT, y, CONTENT_WIDTH, Inches(0.35))
            txbox.text_frame.word_wrap = True
            text = f"[{src.get('id', '')}] {src.get('content', '')}"
            set_text(txbox.text_frame, text,
                     FontStyle(Pt(8), color=colors.secondary_text))
            y += Inches(0.35)


def _slide_disclaimer(new_slide, ctx: dict, colors: PptxColors, styles: dict) -> None:
    """Slide 19: Disclaimer & Attribution."""
    slide = new_slide()
    add_slide_title(slide, "Disclaimer & Attribution", styles, colors)

    disclaimer_text = (
        "This report was generated using AI-powered analysis tools and is provided "
        "for informational purposes only. It does not constitute investment advice, "
        "a recommendation, or an offer to buy or sell any security. The information "
        "contained herein may not be accurate, complete, or current. Investors should "
        "conduct their own research and consult with a qualified financial advisor "
        "before making investment decisions."
    )

    txbox = slide.shapes.add_textbox(CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, Inches(1.5))
    txbox.text_frame.word_wrap = True
    set_text(txbox.text_frame, disclaimer_text, styles["disclaimer"])

    # Attribution block
    y = Inches(3.5)
    analyst_name = ctx.get("analyst_name", "")
    if analyst_name:
        txbox = slide.shapes.add_textbox(CONTENT_LEFT, y, CONTENT_WIDTH, Inches(0.4))
        set_text(txbox.text_frame, analyst_name,
                 FontStyle(Pt(16), bold=True, color=colors.brand), PP_ALIGN.CENTER)
        y += Inches(0.45)

    analyst_info_parts = [p for p in [ctx.get("analyst_sector"), ctx.get("analyst_email")] if p]
    if analyst_info_parts:
        txbox = slide.shapes.add_textbox(CONTENT_LEFT, y, CONTENT_WIDTH, Inches(0.3))
        set_text(txbox.text_frame, " | ".join(analyst_info_parts),
                 FontStyle(Pt(10), color=colors.secondary_text), PP_ALIGN.CENTER)
        y += Inches(0.5)

    firm_name = ctx.get("firm_name", "")
    if firm_name:
        txbox = slide.shapes.add_textbox(CONTENT_LEFT, y, CONTENT_WIDTH, Inches(0.5))
        set_text(txbox.text_frame, firm_name,
                 FontStyle(Pt(18), bold=True, color=colors.brand), PP_ALIGN.CENTER)
        y += Inches(0.6)

    gen_time = ctx.get("generation_time", "")
    tool = ctx.get("tool_branding", "")
    if gen_time or tool:
        footer_parts = [p for p in [tool, gen_time] if p]
        txbox = slide.shapes.add_textbox(CONTENT_LEFT, y, CONTENT_WIDTH, Inches(0.3))
        set_text(txbox.text_frame, " | ".join(footer_parts),
                 FontStyle(Pt(8), color=colors.secondary_text), PP_ALIGN.CENTER)


# =============================================================================
# FOOTBALL FIELD CHART (native PowerPoint shapes)
# =============================================================================

# Impact of changing: alters the vertical density and horizontal scale of
# the football-field chart on the valuation slide.
_FF_LABEL_WIDTH = Inches(2.5)      # space reserved for method labels
_FF_RANGE_LABEL_WIDTH = Inches(1.6)  # space for "$low - $high" text
_FF_BAR_HEIGHT = Inches(0.45)      # height of each range bar
_FF_ROW_GAP = Inches(0.15)        # vertical gap between bars


def _render_football_field_native(
    slide,
    ranges: List[Dict[str, Any]],
    current_price: Optional[float],
    colors: PptxColors,
    styles: dict,
) -> None:
    """Draw the football field as editable PowerPoint shapes.

    Each valuation method gets a horizontal bar from low to high with a
    mid-point diamond marker. A dashed vertical line marks current price.
    """
    n = len(ranges)
    chart_left = CONTENT_LEFT + _FF_LABEL_WIDTH
    chart_width = CONTENT_WIDTH - _FF_LABEL_WIDTH - _FF_RANGE_LABEL_WIDTH
    row_height = _FF_BAR_HEIGHT + _FF_ROW_GAP

    # Compute global price range for scaling
    all_lows = [_safe_float(r.get("low", 0), 0.0) for r in ranges]
    all_highs = [_safe_float(r.get("high", 0), 0.0) for r in ranges]
    price_min = min(v for v in all_lows if v is not None) * 0.9
    price_max = max(v for v in all_highs if v is not None) * 1.1
    price_span = price_max - price_min if price_max > price_min else 1.0

    def price_to_x(price: float) -> int:
        """Convert a price value to an EMU x-position on the chart."""
        frac = (price - price_min) / price_span
        return int(chart_left + chart_width * frac)

    y = CONTENT_TOP

    for r in ranges:
        method = r.get("method", "")
        low = _safe_float(r.get("low", 0), 0.0)
        high = _safe_float(r.get("high", 0), 0.0)
        mid = _safe_float(r.get("mid"))

        # Method label (left-aligned)
        txbox = slide.shapes.add_textbox(CONTENT_LEFT, y, _FF_LABEL_WIDTH - Inches(0.15),
                                         _FF_BAR_HEIGHT)
        txbox.text_frame.word_wrap = True
        set_text(txbox.text_frame, method, styles["body_small"], PP_ALIGN.RIGHT)

        # Range bar background (light gray track)
        track = slide.shapes.add_shape(
            5, chart_left, y + Inches(0.05),
            chart_width, _FF_BAR_HEIGHT - Inches(0.1),
        )
        track.fill.solid()
        track.fill.fore_color.rgb = colors.section_bg
        track.line.fill.background()
        track.adjustments[0] = 0.5  # fully rounded ends

        # Range bar (low to high, brand color with transparency via lighter shade)
        bar_left = price_to_x(low)
        bar_right = price_to_x(high)
        bar_w = max(bar_right - bar_left, Inches(0.1))
        bar = slide.shapes.add_shape(
            5, bar_left, y + Inches(0.05),
            bar_w, _FF_BAR_HEIGHT - Inches(0.1),
        )
        apply_gradient_fill(bar, _lighten(colors.brand, 0.6), _lighten(colors.brand, 0.3),
                            angle_degrees=0.0)
        bar.line.color.rgb = colors.brand
        bar.line.width = Pt(0.75)
        bar.adjustments[0] = 0.5  # rounded ends

        # Mid-point diamond marker
        if mid is not None:
            mid_x = price_to_x(mid)
            diamond_size = Inches(0.2)
            diamond = slide.shapes.add_shape(
                4, mid_x - diamond_size // 2,  # 4 = MSO_SHAPE.DIAMOND
                y + (_FF_BAR_HEIGHT - diamond_size) // 2,
                diamond_size, diamond_size,
            )
            diamond.fill.solid()
            diamond.fill.fore_color.rgb = colors.brand
            diamond.line.fill.background()

        # Range label (right side: "$low - $high")
        label_left = chart_left + chart_width + Inches(0.1)
        txbox = slide.shapes.add_textbox(label_left, y, _FF_RANGE_LABEL_WIDTH, _FF_BAR_HEIGHT)
        set_text(txbox.text_frame,
                 f"${low:,.0f} - ${high:,.0f}" if low and high else "",
                 FontStyle(Pt(8), color=colors.secondary_text))

        y += row_height

    # Current price dashed vertical line
    if current_price is not None:
        cp_x = price_to_x(current_price)
        line_top = CONTENT_TOP
        line_h = int(n * row_height)

        # Draw dashed line as a series of short segments (PowerPoint freeform
        # lines with dashes are fragile; small rectangles are reliable)
        dash_len = Inches(0.12)
        gap_len = Inches(0.08)
        dash_y = line_top
        while dash_y < line_top + line_h:
            seg = slide.shapes.add_shape(1, cp_x, dash_y, Pt(1.5), dash_len)
            seg.fill.solid()
            seg.fill.fore_color.rgb = hex_to_pptx("DC2626")
            seg.line.fill.background()
            dash_y += dash_len + gap_len

        # Current price label
        txbox = slide.shapes.add_textbox(
            cp_x - Inches(1), line_top + line_h + Inches(0.05),
            Inches(2), Inches(0.3),
        )
        set_text(txbox.text_frame, f"Current: ${current_price:,.2f}",
                 FontStyle(Pt(9), bold=True, color=hex_to_pptx("DC2626")),
                 PP_ALIGN.CENTER)


# =============================================================================
# FORMATTING HELPERS
# =============================================================================

def _safe_float(val, default: float = 0.0) -> Optional[float]:
    """Coerce a value to float, returning default on failure."""
    if val is None:
        return None
    try:
        return float(val)
    except (ValueError, TypeError):
        return default


def _fmt_price(val) -> str:
    """Format a price value as $X,XXX.XX or 'N/A'."""
    if val is None:
        return "N/A"
    try:
        v = float(val)
        return f"${v:,.2f}"
    except (ValueError, TypeError):
        return str(val)


def _fmt_pct(val) -> str:
    """Format a percentage value like '+12.3%' or '-5.1%'."""
    if val is None:
        return "N/A"
    try:
        v = float(val)
        sign = "+" if v > 0 else ""
        return f"{sign}{v:.1f}%"
    except (ValueError, TypeError):
        return str(val)


def _fmt_pct_val(val) -> str:
    """Format a raw percentage value (already in % form) like '12.3%'."""
    if val is None:
        return "N/A"
    try:
        v = float(val)
        return f"{v:.1f}%"
    except (ValueError, TypeError):
        return str(val)


def _fmt_number(val) -> str:
    """Format a numeric value with 1 decimal place."""
    if val is None:
        return "N/A"
    try:
        v = float(val)
        return f"{v:,.1f}"
    except (ValueError, TypeError):
        return str(val)


def _fmt_large_number(val) -> str:
    """Format a large number with B/M suffix."""
    if val is None:
        return "N/A"
    try:
        v = float(val)
        if abs(v) >= 1e9:
            return f"${v / 1e9:,.1f}B"
        if abs(v) >= 1e6:
            return f"${v / 1e6:,.1f}M"
        return f"${v:,.0f}"
    except (ValueError, TypeError):
        return str(val)


def _fmt_assumption(val) -> str:
    """Format an assumption value (could be pct or number)."""
    if val is None:
        return "N/A"
    try:
        v = float(val)
        # Heuristic: values < 1 are likely decimals (e.g. 0.08 = 8%)
        if abs(v) < 1:
            return f"{v * 100:.1f}%"
        return f"{v:.1f}%"
    except (ValueError, TypeError):
        return str(val)


def _pct_color(val, colors: PptxColors) -> Optional[Any]:
    """Return green for positive, red for negative, None for zero/None."""
    if val is None:
        return None
    try:
        v = float(val)
        if v > 0:
            return colors.positive
        if v < 0:
            return colors.negative
        return None
    except (ValueError, TypeError):
        return None


def _lighten(color: Any, factor: float) -> Any:
    """Lighten an RGBColor by mixing with white. factor=0.8 means 80% white, 20% color."""
    r = int(color[0] + (255 - color[0]) * factor)
    g = int(color[1] + (255 - color[1]) * factor)
    b = int(color[2] + (255 - color[2]) * factor)
    return type(color)(min(r, 255), min(g, 255), min(b, 255))


def _parse_html_list(html_content: str) -> List[str]:
    """Extract text items from HTML list markup (or plain text lines)."""
    import re
    # Try <li> tags first
    items = re.findall(r'<li[^>]*>(.*?)</li>', html_content, re.DOTALL)
    if items:
        # Strip any remaining HTML tags
        return [re.sub(r'<[^>]+>', '', item).strip() for item in items if item.strip()]
    # Fall back to line splitting
    lines = [line.strip() for line in html_content.split('\n') if line.strip()]
    # Strip bullet characters
    return [re.sub(r'^[\u2022\-\*]\s*', '', line) for line in lines]
