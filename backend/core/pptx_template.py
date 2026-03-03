"""
PowerPoint slide template system.

Provides dimension constants, color presets, font styles, and reusable shape
builders for the PPTX renderer. All measurements are in EMU (English Metric
Units, 914400 per inch) unless noted otherwise.

Impact of changing constants: alters layout of every slide in the generated deck.
Colors are pulled from the shared palette so PPTX stays in sync with PDF/frontend.
"""

from dataclasses import dataclass
from typing import Optional, Sequence, Tuple

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.slide import Slide
from pptx.util import Emu, Inches, Pt

from backend.core.colors import PALETTE
from backend.core.pdf_formatting import strip_markup

# =============================================================================
# SLIDE DIMENSIONS (16:9)
# =============================================================================
# Impact of changing: shifts all content on every slide.
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

MARGIN_LEFT = Inches(0.7)
MARGIN_RIGHT = Inches(0.7)
MARGIN_TOP = Inches(0.95)  # below brand stripe + title
MARGIN_BOTTOM = Inches(0.45)

CONTENT_LEFT = MARGIN_LEFT
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
CONTENT_TOP = Inches(1.05)  # below header band
CONTENT_HEIGHT = SLIDE_HEIGHT - CONTENT_TOP - MARGIN_BOTTOM

# Brand header band at slide top (contains slide title as white text)
STRIPE_HEIGHT = Inches(0.55)

# Footer positioning
FOOTER_TOP = SLIDE_HEIGHT - Inches(0.35)
FOOTER_HEIGHT = Inches(0.25)


# =============================================================================
# COLORS
# =============================================================================
def hex_to_pptx(hex_str: str) -> RGBColor:
    """Convert '#RRGGBB' hex string to pptx RGBColor."""
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


@dataclass(frozen=True)
class PptxColors:
    """Pre-computed RGBColor objects for all semantic uses in slides."""
    brand: RGBColor
    brand_dark: RGBColor
    white: RGBColor
    dark_text: RGBColor
    secondary_text: RGBColor
    label_text: RGBColor
    border_light: RGBColor
    section_bg: RGBColor
    table_alt_row: RGBColor
    positive: RGBColor
    negative: RGBColor
    bull_bg: RGBColor
    bull_border: RGBColor
    bear_bg: RGBColor
    bear_border: RGBColor
    base_case_bg: RGBColor
    estimate_row_bg: RGBColor
    swot_strengths_border: RGBColor
    swot_strengths_bg: RGBColor
    swot_weaknesses_border: RGBColor
    swot_weaknesses_bg: RGBColor
    swot_opportunities_border: RGBColor
    swot_opportunities_bg: RGBColor
    swot_threats_border: RGBColor
    swot_threats_bg: RGBColor


def build_colors(primary_color: Optional[str] = None) -> PptxColors:
    """Build the color set, optionally overriding brand primary from branding config."""
    brand_hex = primary_color or PALETTE.brand.primary
    pdf = PALETTE.pdf
    return PptxColors(
        brand=hex_to_pptx(brand_hex),
        brand_dark=hex_to_pptx(PALETTE.brand.primary_dark),
        white=RGBColor(0xFF, 0xFF, 0xFF),
        dark_text=hex_to_pptx(pdf.dark_text),
        secondary_text=hex_to_pptx(pdf.secondary_text),
        label_text=hex_to_pptx(pdf.label_text),
        border_light=hex_to_pptx(pdf.border_light),
        section_bg=hex_to_pptx(pdf.section_bg),
        table_alt_row=hex_to_pptx(pdf.table_alt_row),
        positive=hex_to_pptx(pdf.positive),
        negative=hex_to_pptx(pdf.negative),
        bull_bg=hex_to_pptx(pdf.bull_case.bg),
        bull_border=hex_to_pptx(pdf.bull_case.border),
        bear_bg=hex_to_pptx(pdf.bear_case.bg),
        bear_border=hex_to_pptx(pdf.bear_case.border),
        base_case_bg=hex_to_pptx(pdf.base_case_bg),
        estimate_row_bg=hex_to_pptx(pdf.estimate_row_bg),
        swot_strengths_border=hex_to_pptx(pdf.swot_strengths.border),
        swot_strengths_bg=hex_to_pptx(pdf.swot_strengths.bg),
        swot_weaknesses_border=hex_to_pptx(pdf.swot_weaknesses.border),
        swot_weaknesses_bg=hex_to_pptx(pdf.swot_weaknesses.bg),
        swot_opportunities_border=hex_to_pptx(pdf.swot_opportunities.border),
        swot_opportunities_bg=hex_to_pptx(pdf.swot_opportunities.bg),
        swot_threats_border=hex_to_pptx(pdf.swot_threats.border),
        swot_threats_bg=hex_to_pptx(pdf.swot_threats.bg),
    )


# =============================================================================
# FONT STYLE PRESETS
# =============================================================================
FONT_NAME = "Calibri"


@dataclass(frozen=True)
class FontStyle:
    """Reusable font specification."""
    size: Pt
    bold: bool = False
    color: Optional[RGBColor] = None
    italic: bool = False


# Applied after colors are built; use factory function instead of module-level.
def font_styles(colors: PptxColors) -> dict:
    """Return a dict of named FontStyle presets keyed to the given color set."""
    return {
        "slide_title": FontStyle(Pt(16), bold=True, color=colors.white),
        "body": FontStyle(Pt(10), color=colors.dark_text),
        "body_small": FontStyle(Pt(9), color=colors.dark_text),
        "metric_value": FontStyle(Pt(18), bold=True, color=colors.brand),
        "metric_value_sm": FontStyle(Pt(14), bold=True, color=colors.brand),
        "metric_label": FontStyle(Pt(8), color=colors.secondary_text),
        "table_header": FontStyle(Pt(8), bold=True, color=colors.white),
        "table_cell": FontStyle(Pt(8), color=colors.dark_text),
        "footer": FontStyle(Pt(7), color=colors.secondary_text),
        "bullet": FontStyle(Pt(9), color=colors.dark_text),
        "section_label": FontStyle(Pt(10), bold=True, color=colors.label_text),
        "disclaimer": FontStyle(Pt(8), color=colors.secondary_text, italic=True),
    }


# =============================================================================
# TEXT HELPERS
# =============================================================================
def apply_font(run, style: FontStyle) -> None:
    """Apply a FontStyle to a pptx Run object."""
    run.font.name = FONT_NAME
    run.font.size = style.size
    run.font.bold = style.bold
    run.font.italic = style.italic
    if style.color is not None:
        run.font.color.rgb = style.color


def set_text(text_frame, text: str, style: FontStyle, alignment: int = PP_ALIGN.LEFT) -> None:
    """Set text on a text frame's first paragraph with the given style."""
    text_frame.clear()
    text_frame.word_wrap = True
    para = text_frame.paragraphs[0]
    para.alignment = alignment
    run = para.add_run()
    run.text = str(text)
    apply_font(run, style)


def add_paragraph(text_frame, text: str, style: FontStyle, alignment: int = PP_ALIGN.LEFT):
    """Append a new paragraph to an existing text frame."""
    para = text_frame.add_paragraph()
    para.alignment = alignment
    para.space_before = Pt(3)
    para.space_after = Pt(1)
    run = para.add_run()
    run.text = str(text)
    apply_font(run, style)
    return para


# =============================================================================
# SHAPE BUILDERS
# =============================================================================
def add_brand_stripe(slide: Slide, colors: PptxColors) -> None:
    """Add a thin colored rectangle across the top of the slide."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Emu(0), Emu(0), SLIDE_WIDTH, STRIPE_HEIGHT,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors.brand
    shape.line.fill.background()


def add_footer(slide: Slide, firm_name: str, slide_num: int, colors: PptxColors, styles: dict) -> None:
    """Add a separator line above the footer, firm name (left) and slide number (right)."""
    # Thin separator line above footer
    sep_top = FOOTER_TOP - Inches(0.08)
    sep = slide.shapes.add_shape(1, CONTENT_LEFT, sep_top, CONTENT_WIDTH, Pt(0.5))
    sep.fill.solid()
    sep.fill.fore_color.rgb = colors.border_light
    sep.line.fill.background()

    # Firm name
    txbox = slide.shapes.add_textbox(MARGIN_LEFT, FOOTER_TOP, Inches(4), FOOTER_HEIGHT)
    set_text(txbox.text_frame, firm_name, styles["footer"])

    # Slide number
    txbox = slide.shapes.add_textbox(
        SLIDE_WIDTH - MARGIN_RIGHT - Inches(1), FOOTER_TOP, Inches(1), FOOTER_HEIGHT,
    )
    set_text(txbox.text_frame, str(slide_num), styles["footer"], PP_ALIGN.RIGHT)


def add_slide_title(slide: Slide, title: str, styles: dict, colors: Optional["PptxColors"] = None) -> None:
    """Place the slide title as white text inside the brand header band.

    The title text is vertically centered within the stripe shape. Word wrap
    is enabled to handle long titles (tested up to ~80 characters).
    """
    txbox = slide.shapes.add_textbox(
        MARGIN_LEFT, Emu(0), CONTENT_WIDTH, STRIPE_HEIGHT,
    )
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    # Vertically center text within the header band
    body_pr = tf._txBody.find(qn("a:bodyPr"))
    if body_pr is not None:
        body_pr.set("anchor", "ctr")
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT
    run = para.add_run()
    run.text = str(title)
    apply_font(run, styles["slide_title"])


def add_metric_card(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    label: str,
    value: str,
    colors: PptxColors,
    styles: dict,
    value_color: Optional[RGBColor] = None,
) -> None:
    """Add a metric card: rounded rectangle with label above value.

    Padding is generous (matching the PDF's 5mm/6mm feel) so that values
    breathe rather than butting up against the edges.
    """
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE (sharp corners for institutional look)
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors.section_bg
    shape.line.color.rgb = colors.border_light
    shape.line.width = Pt(0.5)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.1)

    # Label
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    para.space_after = Pt(4)
    run = para.add_run()
    run.text = str(label).upper()
    apply_font(run, styles["metric_label"])

    # Value
    para = tf.add_paragraph()
    para.alignment = PP_ALIGN.CENTER
    para.space_before = Pt(4)
    run = para.add_run()
    run.text = str(value)
    v_style = FontStyle(
        size=styles["metric_value"].size,
        bold=True,
        color=value_color or styles["metric_value"].color,
    )
    apply_font(run, v_style)


def add_bullet_card(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    accent_color: RGBColor,
    colors: PptxColors,
    styles: dict,
) -> int:
    """Add a bullet card with a colored left accent bar. Returns the bottom edge position.

    The accent bar is slightly thicker than the original (matching the PDF's
    2.5pt border-left) and the text has more left padding for breathing room.
    """
    # Strip markdown/HTML that LLM-generated text may contain
    text = strip_markup(text)

    # Accent bar (thicker, matching PDF's 2.5pt border-left)
    bar_width = Inches(0.07)
    bar = slide.shapes.add_shape(1, left, top, bar_width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()

    # Text box with more generous left margin
    text_left = left + bar_width + Inches(0.15)
    text_width = width - bar_width - Inches(0.15)
    txbox = slide.shapes.add_textbox(text_left, top, text_width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.05)
    set_text(tf, text, styles["bullet"])

    return top + height


def add_bullet_list(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    bullets: Sequence[str],
    accent_color: RGBColor,
    colors: PptxColors,
    styles: dict,
    card_height: Optional[int] = None,
    spacing: Optional[int] = None,
) -> int:
    """Add a vertical stack of bullet cards. Returns the bottom edge position."""
    if not bullets:
        return top
    per_card = card_height or Inches(0.75)
    gap = spacing or Inches(0.12)
    y = top
    for bullet in bullets:
        y = add_bullet_card(slide, left, y, width, per_card, bullet, accent_color, colors, styles)
        y += gap
    return y


def create_table(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    row_count: int,
    col_count: int,
    col_widths: Optional[Sequence[int]] = None,
) -> "Table":
    """Create a native PowerPoint table and return the Table object.

    If col_widths is provided, it should be a sequence of EMU values for each column.
    Otherwise columns are distributed evenly.
    """
    row_height = Inches(0.35)
    height = row_height * row_count
    graphic_frame = slide.shapes.add_table(row_count, col_count, left, top, width, height)
    table = graphic_frame.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    else:
        even_width = width // col_count
        for i in range(col_count):
            table.columns[i].width = even_width

    return table


def style_table(
    table,
    colors: PptxColors,
    styles: dict,
    header_labels: Optional[Sequence[str]] = None,
) -> None:
    """Apply branded styling to a table: blue header row, alternating body rows.

    Uses light dotted bottom borders on data cells instead of PowerPoint's
    default heavy solid borders, matching the PDF template's refined look.
    """
    # Disable PowerPoint's built-in table style banding to avoid conflicts
    _disable_table_style(table)

    border_hex = _rgb_to_hex(colors.border_light)

    for row_idx, row in enumerate(table.rows):
        for col_idx, cell in enumerate(row.cells):
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = colors.brand
                if header_labels and col_idx < len(header_labels):
                    set_text(cell.text_frame, header_labels[col_idx], styles["table_header"])
                # Clean borders on header: none on sides, solid bottom
                _set_cell_borders(cell, bottom_w=Pt(1), bottom_color=_rgb_to_hex(colors.brand))
            else:
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = hex_to_pptx("F7F8FA")
                else:
                    cell.fill.background()
                # Light dotted bottom border, no side borders
                _set_cell_borders(
                    cell,
                    bottom_w=Pt(0.5),
                    bottom_color=border_hex,
                    bottom_dash="sysDot",
                )


def set_cell(cell, text: str, style: FontStyle, alignment: int = PP_ALIGN.LEFT) -> None:
    """Set text on a table cell."""
    set_text(cell.text_frame, text, style, alignment)


# =============================================================================
# TABLE BORDER HELPERS (XML-level, since python-pptx lacks a high-level API)
# =============================================================================

# Impact of changing: affects border rendering on every table cell in the deck.
_BORDER_SIDES = ("a:lnL", "a:lnR", "a:lnT", "a:lnB")


def _set_cell_borders(
    cell,
    *,
    bottom_w: int = 0,
    bottom_color: str = "D1D5DB",
    bottom_dash: str = "solid",
) -> None:
    """Set cell borders via XML. Suppresses all borders except a styled bottom.

    bottom_w     -- border width as EMU (use Pt() to convert from points)
    bottom_color -- six-char hex RGB, e.g. 'D1D5DB'
    bottom_dash  -- OOXML dash preset: 'solid', 'sysDot', 'dot', 'dash'
    """
    tc_pr = cell._tc.get_or_add_tcPr()

    # Remove any existing border elements to avoid PowerPoint repair warnings
    for tag in _BORDER_SIDES:
        for existing in tc_pr.findall(qn(tag)):
            tc_pr.remove(existing)

    # Suppress left, right, top borders
    for tag in ("a:lnL", "a:lnR", "a:lnT"):
        ln = _make_no_border(tag)
        tc_pr.append(ln)

    # Bottom border: styled or suppressed
    if bottom_w and bottom_w > 0:
        tc_pr.append(_make_border("a:lnB", int(bottom_w), bottom_color, bottom_dash))
    else:
        tc_pr.append(_make_no_border("a:lnB"))


def _make_border(tag: str, width_emu: int, color_hex: str, dash: str) -> etree._Element:
    """Build an <a:lnL/R/T/B> border element with specified style."""
    ln = etree.Element(qn(tag))
    ln.set("w", str(width_emu))
    ln.set("cap", "flat")
    ln.set("cmpd", "sng")

    fill = etree.SubElement(ln, qn("a:solidFill"))
    clr = etree.SubElement(fill, qn("a:srgbClr"))
    clr.set("val", color_hex)

    prd = etree.SubElement(ln, qn("a:prstDash"))
    prd.set("val", dash)
    return ln


def _make_no_border(tag: str) -> etree._Element:
    """Build a zero-width border element that suppresses the line."""
    ln = etree.Element(qn(tag))
    ln.set("w", "0")
    etree.SubElement(ln, qn("a:noFill"))
    return ln


def _disable_table_style(table) -> None:
    """Clear the built-in table style so our custom borders aren't overridden.

    PowerPoint tables inherit banding and border styles from a:tblStyleId.
    Setting it to the "no style" GUID lets our per-cell formatting take effect.
    """
    tbl = table._tbl
    tbl_pr = tbl.tblPr
    if tbl_pr is not None:
        # {2D5ABB26-0587-4C30-8999-92F81FD0307C} = No Style, No Grid
        for style_id in tbl_pr.findall(qn("a:tblStyleId")):
            tbl_pr.remove(style_id)
        # Also disable first-row / banding flags so they don't inject borders
        tbl_pr.set("firstRow", "0")
        tbl_pr.set("bandRow", "0")
        tbl_pr.set("lastRow", "0")
        tbl_pr.set("firstCol", "0")
        tbl_pr.set("lastCol", "0")


def _rgb_to_hex(color: RGBColor) -> str:
    """Convert an RGBColor to a six-char hex string (no '#' prefix)."""
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


# =============================================================================
# CENTERING HELPER
# =============================================================================

def centered_left(n_items: int, item_width: int, gap: int) -> int:
    """Return the left position that centers a row of equally-sized items within CONTENT_WIDTH.

    All values are in EMU. The row total is n_items * item_width + (n_items - 1) * gap,
    and the returned position is the left edge of the first item.
    """
    total = n_items * item_width + max(0, n_items - 1) * gap
    return int(CONTENT_LEFT + (CONTENT_WIDTH - total) // 2)


# =============================================================================
# GRADIENT FILL HELPER
# =============================================================================

def apply_gradient_fill(
    shape,
    color_start: RGBColor,
    color_end: RGBColor,
    angle_degrees: float = 0.0,
) -> None:
    """Apply a two-stop linear gradient fill to a shape.

    angle_degrees -- 0 = left-to-right, 90 = bottom-to-top, 270 = top-to-bottom.
    Uses the python-pptx high-level gradient API.
    """
    shape.fill.gradient()
    stops = shape.fill.gradient_stops
    stops[0].position = 0.0
    stops[0].color.rgb = color_start
    stops[1].position = 1.0
    stops[1].color.rgb = color_end
    shape.fill.gradient_angle = angle_degrees
