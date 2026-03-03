"""Tests for PowerPoint slide deck generation.

Validates that generate_pptx produces valid PPTX bytes with the expected
slide structure, and that individual formatting helpers behave correctly.
"""

import io

import pytest
from pptx import Presentation

from backend.core.pptx_renderer import (
    _fmt_large_number,
    _fmt_pct,
    _fmt_price,
    _lighten,
    _parse_html_list,
    generate_pptx,
)
from backend.core.pptx_template import PptxColors, build_colors, font_styles, hex_to_pptx


# =============================================================================
# FIXTURES
# =============================================================================

@pytest.fixture()
def minimal_context() -> dict:
    """Context with only required fields: produces title + disclaimer slides."""
    return {
        "primary_color": "#1E3A5F",
        "firm_name": "Test Firm",
        "tool_branding": "Platform",
        "analyst_name": "Analyst",
        "analyst_email": "a@b.com",
        "analyst_sector": "Tech",
        "ticker": "TEST",
        "company_name": "Test Corp",
        "exchange": "NYSE",
        "sector": "Technology",
        "rating": "BUY",
        "rating_color": "#16a34a",
        "report_date": "01 January 2026",
        "generation_time": "2026-01-01 00:00",
    }


@pytest.fixture()
def full_context(minimal_context) -> dict:
    """Context with all optional fields populated."""
    return {
        **minimal_context,
        "toc_entries": [{"title": "Investment Thesis"}, {"title": "DCF Valuation"}],
        "exec_summary": {
            "rating": "BUY",
            "target_price": 150.0,
            "current_price": 100.0,
            "upside_pct": 50.0,
            "catalysts": ["Catalyst A", "Catalyst B"],
            "key_risk": "Key risk here",
            "weighted_fv": 148.0,
        },
        "report_headline": "Strong buy on fundamentals",
        "thesis_bullets": ["Point 1", "Point 2", "Point 3"],
        "bull_bullets": ["Bull 1", "Bull 2"],
        "bear_bullets": ["Bear 1"],
        "moat_bullets": ["Moat 1"],
        "external_bullets": ["Industry 1"],
        "conviction_structured": {"overall_score": 82, "confidence": 90},
        "conviction_categories": [
            {"name": "Fundamentals", "score": 8.5, "bar_pct": 85, "evidence": "Strong"},
            {"name": "Valuation", "score": 6.0, "bar_pct": 60, "evidence": "Fair"},
        ],
        "conviction_summary": "Conviction is high. Valuation supports thesis.",
        "key_data": [
            {"label": "Price", "value": "$100.00"},
            {"label": "Market Cap", "value": "$50B"},
        ],
        "statistics": [{"label": "P/E", "value": "25.0x"}],
        "returns": [
            {"period": "1M", "value": 3.5, "display": "+3.5%"},
            {"period": "YTD", "value": -2.1, "display": "-2.1%"},
        ],
        "dcf_summary": {
            "fair_value": 148.0,
            "upside_pct": 48.0,
            "wacc": 9.0,
            "terminal_growth_rate": 2.5,
            "enterprise_value": 60_000_000_000,
            "tv_pct_of_ev": 62.0,
        },
        "scenario_summary": {
            "bull": {"fair_value": 200, "probability": 25, "upside_pct": 100, "assumptions": {"revenue_growth": 0.15}},
            "base": {"fair_value": 148, "probability": 55, "upside_pct": 48, "assumptions": {}},
            "bear": {"fair_value": 80, "probability": 20, "upside_pct": -20, "assumptions": {}},
            "weighted_fair_value": 150,
        },
        "football_ranges": [
            {"method": "DCF", "low": 130, "mid": 148, "high": 170},
            {"method": "P/E Comps", "low": 120, "mid": 140, "high": 160},
        ],
        "football_current_pct": 40,
        "football_current_price": 100.0,
        "sensitivity_grid": {
            "row_values": [8.0, 9.0, 10.0],
            "col_values": [2.0, 2.5, 3.0],
            "grid": [[180, 165, 155], [160, 148, 140], [145, 135, 128]],
            "cell_classes": [
                ["positive", "positive", "positive"],
                ["positive", "base", "positive"],
                ["positive", "positive", "positive"],
            ],
            "base_row_idx": 1,
            "base_col_idx": 1,
        },
        "growth_margin_grid": None,
        "earnings_table": {
            "rows": [
                {"fiscal_year": "FY2024", "revenue": 50_000_000_000, "revenue_growth": 10.0,
                 "ebitda": 15_000_000_000, "ebitda_margin": 30.0, "eps": 5.50, "is_estimate": False},
                {"fiscal_year": "FY2025E", "revenue": 56_000_000_000, "revenue_growth": 12.0,
                 "ebitda": 17_500_000_000, "ebitda_margin": 31.3, "eps": 6.30, "is_estimate": True},
            ],
            "has_estimates": True,
        },
        "earnings_takeaway": "Revenue growth accelerating with margin expansion.",
        "comp_summary": {
            "target": {"ticker": "TEST", "pe_ratio": 25.0, "forward_pe": 22.0,
                       "ev_to_ebitda": 18.0, "ev_to_revenue": 5.0, "market_cap": 50_000_000_000,
                       "gross_margin": 55.0, "revenue_growth": 10.0},
            "peers": [
                {"ticker": "PEER1", "pe_ratio": 30.0, "forward_pe": 26.0,
                 "ev_to_ebitda": 20.0, "ev_to_revenue": 6.0, "market_cap": 40_000_000_000,
                 "gross_margin": 50.0, "revenue_growth": 8.0},
            ],
            "medians": {"pe_ratio": 27.5, "forward_pe": 24.0, "ev_to_ebitda": 19.0,
                        "ev_to_revenue": 5.5, "gross_margin": 52.5, "revenue_growth": 9.0},
            "implied_values": {"pe_implied": 155, "ev_ebitda_implied": 160, "ev_revenue_implied": 145},
        },
        "precedent_summary": {
            "deals": [
                {"date": "2023-06", "acquirer": "AcqCo", "target": "TgtCo",
                 "ev_to_revenue": 6.0, "ev_to_ebitda": 20.0, "premium_pct": 40.0},
            ],
            "medians": {"ev_revenue": 6.0, "ev_ebitda": 20.0, "premium": 40.0},
            "implied_values": {"ev_revenue": 160, "ev_ebitda": 170},
        },
        "precedent_takeaway": "Median deal premium of 40% vs current upside of 48%.",
        "strategy_grid": {
            "advantages": "<li>Strong brand</li><li>Ecosystem</li>",
            "vulnerabilities": "<li>Concentration risk</li>",
            "tailwinds": "<li>Market growth</li>",
            "headwinds": "<li>Regulation</li>",
        },
        "strategy_bullets": [],
        "research_sources": [
            {"id": 1, "title": "Annual Report FY2024", "url": "https://example.com", "date": "2024-11"},
        ],
        "analyst_sources": [{"id": "A1", "content": "Management guidance notes"}],
        "model_audit": {},
    }


# =============================================================================
# TEMPLATE MODULE TESTS
# =============================================================================

class TestPptxTemplate:
    def test_hex_to_pptx_valid(self):
        color = hex_to_pptx("#1E3A5F")
        assert color[0] == 0x1E
        assert color[1] == 0x3A
        assert color[2] == 0x5F

    def test_hex_to_pptx_no_hash(self):
        color = hex_to_pptx("FF0000")
        assert color[0] == 0xFF
        assert color[1] == 0x00
        assert color[2] == 0x00

    def test_build_colors_defaults(self):
        colors = build_colors()
        assert isinstance(colors, PptxColors)
        assert colors.white[0] == 0xFF

    def test_build_colors_custom_primary(self):
        colors = build_colors("#FF0000")
        assert colors.brand[0] == 0xFF
        assert colors.brand[1] == 0x00

    def test_font_styles_all_present(self):
        colors = build_colors()
        styles = font_styles(colors)
        expected = {"slide_title", "body", "body_small", "metric_value", "metric_value_sm",
                    "metric_label", "table_header", "table_cell", "footer", "bullet",
                    "section_label", "disclaimer"}
        assert set(styles.keys()) == expected


# =============================================================================
# FORMATTING HELPERS
# =============================================================================

class TestFormatters:
    @pytest.mark.parametrize("val,expected", [
        (100.0, "$100.00"),
        (1234.56, "$1,234.56"),
        (None, "N/A"),
        ("invalid", "invalid"),
    ])
    def test_fmt_price(self, val, expected):
        assert _fmt_price(val) == expected

    @pytest.mark.parametrize("val,expected", [
        (25.0, "+25.0%"),
        (-10.5, "-10.5%"),
        (0, "0.0%"),
        (None, "N/A"),
    ])
    def test_fmt_pct(self, val, expected):
        assert _fmt_pct(val) == expected

    @pytest.mark.parametrize("val,expected", [
        (3_200_000_000_000, "$3,200.0B"),
        (50_000_000, "$50.0M"),
        (1234, "$1,234"),
        (None, "N/A"),
    ])
    def test_fmt_large_number(self, val, expected):
        assert _fmt_large_number(val) == expected

    def test_lighten(self):
        from pptx.dml.color import RGBColor
        color = RGBColor(0x1E, 0x3A, 0x5F)
        lightened = _lighten(color, 0.5)
        # Should be closer to white
        assert lightened[0] > color[0]
        assert lightened[1] > color[1]
        assert lightened[2] > color[2]

    @pytest.mark.parametrize("html,expected_count", [
        ("<li>Item 1</li><li>Item 2</li>", 2),
        ("Line 1\nLine 2\nLine 3", 3),
        ("- Bullet 1\n- Bullet 2", 2),
        ("", 0),
    ])
    def test_parse_html_list(self, html, expected_count):
        result = _parse_html_list(html)
        assert len(result) == expected_count


# =============================================================================
# FULL GENERATION TESTS
# =============================================================================

class TestGeneratePptx:
    def test_minimal_context_produces_valid_pptx(self, minimal_context):
        """Minimal context should produce at least title + disclaimer slides."""
        result = generate_pptx(minimal_context)
        assert isinstance(result, bytes)
        assert len(result) > 0

        prs = Presentation(io.BytesIO(result))
        # Title slide + disclaimer at minimum
        assert len(prs.slides) >= 2

    def test_full_context_produces_all_slides(self, full_context):
        """Full context should produce all 19 slides."""
        result = generate_pptx(full_context)
        prs = Presentation(io.BytesIO(result))
        assert len(prs.slides) == 19

    def test_slide_dimensions_are_widescreen(self, full_context):
        """Slides should be 16:9 widescreen format."""
        result = generate_pptx(full_context)
        prs = Presentation(io.BytesIO(result))
        # 16:9 = 13.333" x 7.5" in EMU (Inches() truncates, so use tolerance)
        assert abs(prs.slide_width - 12192000) < 1000
        assert abs(prs.slide_height - 6858000) < 1000

    def test_title_slide_has_company_name(self, full_context):
        """First slide should contain the company name."""
        result = generate_pptx(full_context)
        prs = Presentation(io.BytesIO(result))
        slide = prs.slides[0]
        all_text = _extract_slide_text(slide)
        assert "Test Corp" in all_text

    def test_title_slide_has_rating(self, full_context):
        """Title slide should show the rating."""
        result = generate_pptx(full_context)
        prs = Presentation(io.BytesIO(result))
        slide = prs.slides[0]
        all_text = _extract_slide_text(slide)
        assert "BUY" in all_text

    def test_missing_section_skips_slide(self, minimal_context):
        """Missing optional data should not create a slide for that section."""
        minimal_context["thesis_bullets"] = []
        result = generate_pptx(minimal_context)
        prs = Presentation(io.BytesIO(result))
        all_text = " ".join(_extract_slide_text(s) for s in prs.slides)
        assert "Investment Thesis" not in all_text

    def test_football_field_slide_present(self, full_context):
        """Football field slide should render (either as chart or table)."""
        result = generate_pptx(full_context)
        prs = Presentation(io.BytesIO(result))
        all_titles = [_extract_slide_text(s) for s in prs.slides]
        assert any("Valuation Football Field" in t for t in all_titles)

    def test_pptx_is_reopenable(self, full_context):
        """Generated PPTX should be parseable by python-pptx (valid format)."""
        result = generate_pptx(full_context)
        prs = Presentation(io.BytesIO(result))
        # Iterate all shapes to ensure no corrupt references
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    _ = shape.text_frame.text

    def test_no_empty_slides(self, full_context):
        """Every slide should have at least one shape."""
        result = generate_pptx(full_context)
        prs = Presentation(io.BytesIO(result))
        for i, slide in enumerate(prs.slides):
            assert len(slide.shapes) > 0, f"Slide {i + 1} has no shapes"

    def test_resilient_to_bad_data(self, minimal_context):
        """Bad data in a section should skip that slide, not crash the whole deck."""
        minimal_context["dcf_summary"] = {"fair_value": "not_a_number"}
        result = generate_pptx(minimal_context)
        assert isinstance(result, bytes)
        assert len(result) > 0


# =============================================================================
# HELPERS
# =============================================================================

def _extract_slide_text(slide) -> str:
    """Extract all text from a slide as a single string."""
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    parts.append(t)
    return " ".join(parts)
